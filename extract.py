from pptx import Presentation
import glob
import os

# Define the pattern to match PPTX files in the current directory
pptx_pattern = "*.pptx"

# Use glob to find PPTX files in the current directory
pptx_files = glob.glob(pptx_pattern)

if pptx_files:
    for pptx_file in pptx_files:
        prs = Presentation(pptx_file)

        # Create a text file for each PowerPoint file
        output_file = os.path.splitext(pptx_file)[0] + '.txt'

        with open(output_file, 'w', encoding='utf-8') as file:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            # Save the run.text content to the file
                            file.write(run.text + '\n')
else:
    print("No PPTX files found in the directory.") 
